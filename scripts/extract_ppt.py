#!/usr/bin/env python3
"""
PPT Table Extraction Script

Extracts slide content including tables from PowerPoint files for RAG ingestion.
Outputs structured JSON with slide numbers, titles, body text, and table data.
"""

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table


def extract_table_data(table: Table) -> list[list[str]]:
    """
    Extract table data as a 2D list of strings, preserving rows and columns.

    Args:
        table: A python-pptx Table object

    Returns:
        List of rows, where each row is a list of cell text values
    """
    rows = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            # Get text and preserve numerical precision
            cell_text = cell.text.strip() if cell.text else ""
            row_data.append(cell_text)
        rows.append(row_data)
    return rows


def extract_shape_text(shape: BaseShape) -> str:
    """
    Extract text from a shape if it has a text frame.

    Args:
        shape: A python-pptx shape object

    Returns:
        Extracted text or empty string
    """
    if hasattr(shape, "text_frame"):
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs)
            if text.strip():
                paragraphs.append(text.strip())
        return "\n".join(paragraphs)
    return ""


def extract_slide_content(slide, slide_number: int) -> dict[str, Any]:
    """
    Extract all content from a single slide.

    Args:
        slide: A python-pptx slide object
        slide_number: 1-indexed slide number

    Returns:
        Dictionary containing slide number, title, body text, and tables
    """
    slide_data: dict[str, Any] = {
        "slide_number": slide_number,
        "title": "",
        "body_text": [],
        "tables": []
    }

    for shape in slide.shapes:
        # Extract title
        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
            from pptx.enum.shapes import PP_PLACEHOLDER
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                slide_data["title"] = extract_shape_text(shape)
                continue

        # Extract tables
        if shape.has_table:
            table_data = extract_table_data(shape.table)
            if table_data and any(any(cell for cell in row) for row in table_data):
                slide_data["tables"].append({
                    "rows": len(table_data),
                    "columns": len(table_data[0]) if table_data else 0,
                    "data": table_data
                })
        # Extract text from other shapes
        elif hasattr(shape, "text_frame"):
            text = extract_shape_text(shape)
            if text and text != slide_data["title"]:
                slide_data["body_text"].append(text)

    return slide_data


def extract_ppt(ppt_path: Path) -> dict[str, Any]:
    """
    Extract all content from a PowerPoint file.

    Args:
        ppt_path: Path to the .pptx file

    Returns:
        Dictionary containing file metadata and all slide content
    """
    prs = Presentation(str(ppt_path))

    extraction_result = {
        "source_file": ppt_path.name,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, start=1):
        slide_content = extract_slide_content(slide, idx)
        extraction_result["slides"].append(slide_content)

    # Count tables
    total_tables = sum(len(s["tables"]) for s in extraction_result["slides"])
    extraction_result["total_tables"] = total_tables

    return extraction_result


def process_ppt_files(input_dir: Path, output_dir: Path) -> dict[str, Any]:
    """
    Process all PPTX files in the input directory.

    Args:
        input_dir: Directory containing PPTX files
        output_dir: Directory for output JSON files

    Returns:
        Summary of processing results
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    ppt_files = list(input_dir.glob("*.pptx"))

    # Filter out template files
    ppt_files = [f for f in ppt_files if "template" not in f.name.lower() and "tamplate" not in f.name.lower()]

    summary = {
        "files_processed": [],
        "total_files": 0,
        "total_slides": 0,
        "total_tables": 0,
        "errors": []
    }

    for ppt_path in sorted(ppt_files):
        print(f"Processing: {ppt_path.name}")
        try:
            result = extract_ppt(ppt_path)

            # Generate output filename
            output_filename = ppt_path.stem + ".json"
            output_path = output_dir / output_filename

            # Write JSON with proper encoding for Chinese characters
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(result, f, ensure_ascii=False, indent=2)

            file_summary = {
                "source": ppt_path.name,
                "output": output_filename,
                "slides": result["total_slides"],
                "tables": result["total_tables"]
            }
            summary["files_processed"].append(file_summary)
            summary["total_slides"] += result["total_slides"]
            summary["total_tables"] += result["total_tables"]

            print(f"  - Extracted {result['total_slides']} slides, {result['total_tables']} tables")
            print(f"  - Saved to: {output_path}")

        except Exception as e:
            error_msg = f"Error processing {ppt_path.name}: {str(e)}"
            print(f"  - ERROR: {error_msg}")
            summary["errors"].append(error_msg)

    summary["total_files"] = len(summary["files_processed"])
    return summary


def main():
    """Main entry point."""
    # Default paths
    default_input = Path("/home/guilinzhang/allProjects/mfg-ai-platform/docs/ppt_files/ppt")
    default_output = Path("/home/guilinzhang/allProjects/mfg-ai-platform/data/extracted")

    # Allow command line overrides
    input_dir = Path(sys.argv[1]) if len(sys.argv) > 1 else default_input
    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else default_output

    if not input_dir.exists():
        print(f"Error: Input directory does not exist: {input_dir}")
        sys.exit(1)

    print("=" * 60)
    print("PPT Table Extraction Script")
    print("=" * 60)
    print(f"Input directory:  {input_dir}")
    print(f"Output directory: {output_dir}")
    print("=" * 60)
    print()

    summary = process_ppt_files(input_dir, output_dir)

    print()
    print("=" * 60)
    print("EXTRACTION SUMMARY")
    print("=" * 60)
    print(f"Files processed: {summary['total_files']}")
    print(f"Total slides:    {summary['total_slides']}")
    print(f"Total tables:    {summary['total_tables']}")

    if summary["errors"]:
        print(f"\nErrors encountered: {len(summary['errors'])}")
        for error in summary["errors"]:
            print(f"  - {error}")

    print()
    print("Per-file breakdown:")
    for file_info in summary["files_processed"]:
        print(f"  {file_info['source']}")
        print(f"    -> {file_info['slides']} slides, {file_info['tables']} tables")

    # Save summary
    summary_path = output_dir / "extraction_summary.json"
    with open(summary_path, "w", encoding="utf-8") as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    print(f"\nSummary saved to: {summary_path}")

    return 0 if not summary["errors"] else 1


if __name__ == "__main__":
    sys.exit(main())
